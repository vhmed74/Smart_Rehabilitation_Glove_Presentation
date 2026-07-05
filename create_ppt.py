import os
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(title, content_list, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for i, txt in enumerate(content_list):
        if i == 0:
            tf.text = txt
        else:
            p = tf.add_paragraph()
            p.text = txt
    
    if image_path and os.path.exists(image_path):
        left = Inches(5.5)
        top = Inches(2.0)
        # Adding image and scaling to max 4.5 inches tall
        slide.shapes.add_picture(image_path, left, top, height=Inches(4.5))

# Slide 1: Hero
add_title_slide("Smart Rehabilitation Glove", "Graduation Project 2026\nMenoufia University\nSupervised By: Dr. Essam Nabil")

# Slide 2: Abstract
add_content_slide("Abstract", [
    "An Integrated Therapeutic System.",
    "A smart wearable rehabilitation glove for patients with neurological impairments.",
    "TENS Module: Transcutaneous Electrical Nerve Stimulation.",
    "Vibration Therapy Module: ERM haptic actuators.",
    "Force Measurement Module: FSR sensors mapped to 0-100% grip."
], "images/image9.png")

# Slide 3: Problem Statement
add_content_slide("Problem Statement", [
    "Core Challenges of Hand Rehabilitation:",
    "- Limited Range of Motion",
    "- Muscle Atrophy & Weakness",
    "- Sensory Deficits",
    "- Impaired Coordination",
    "- Chronic Pain",
    "Our Smart Solution: An integrated wearable glove."
])

# Slide 4: Historical Background
add_content_slide("Historical Background", [
    "Evolution of Hand Rehabilitation:",
    "- Ancient Era: Primitive splinting",
    "- 18th-19th Century: Surgical Pioneers",
    "- World Wars: Wartime Innovation & FES",
    "- 1950s-1977: Formal Specialty",
    "- Today 2026: Smart Wearables"
])

# Slide 5: Objectives (1/2)
add_content_slide("Objectives (1/2)", [
    "What We Aim to Achieve:",
    "1. Restoring Range of Motion (ROM)",
    "2. Increasing Strength & Dexterity",
    "3. Non-Invasive Pain Management"
], "images/image7.png")

# Slide 6: Objectives (2/2)
add_content_slide("Objectives (2/2)", [
    "Beyond Basic Recovery:",
    "4. Enhancing Activities of Daily Living",
    "5. Improving Coordination & Proprioception",
    "Customized Plans for targeted therapy.",
    "Scar Management techniques."
])

# Slide 7: System Architecture
add_content_slide("System Architecture", [
    "4-Layer Design Philosophy:",
    "1. Input & Sensing Layer: Flex + FSR sensors",
    "2. Processing & Control: ESP32 dual-core",
    "3. Actuation & Power: Solenoid valves + MOSFET",
    "4. Feedback & Interface: LCD display"
], "images/image6.jpeg")

# Slide 8: Components (Core)
add_content_slide("Materials & Methods: Core Components", [
    "Control, Power & Actuation:",
    "- ESP32 Microcontroller",
    "- Pneumatic Glove",
    "- LCD Display Screen",
    "- MOSFET Module",
    "- Power Supply Adaptor",
    "- Voltage Regulator (LM2596)"
], "images/image9.png")

# Slide 9: Components (Sensors)
add_content_slide("Materials & Methods: Sensors & Therapy", [
    "Sensing, Feedback & Therapy:",
    "- TENS Cable & Electrodes",
    "- Vibration Motor (ERM)",
    "- Force Sensor (FSR)",
    "- Solenoid Valves",
    "- Voltage Regulator 7805",
    "- Air Pressure Sensor HX710B"
], "images/image14.png")

# Slide 10: Dev Methodology (1/2)
add_content_slide("Development Methodology (1/2)", [
    "6-Phase Top-Down Approach:",
    "1. Analysis & Requirements",
    "2. Theoretical Design",
    "3. Modular Assembly & Unit Testing"
], "images/image3.jpeg")

# Slide 11: Dev Methodology (2/2)
add_content_slide("Development Methodology (2/2)", [
    "Integration & Calibration:",
    "4. Full System Integration",
    "5. Functional Verification",
    "6. Calibration & Final Tuning",
    "Developed with Arduino IDE (C++) and CAD software."
])

# Slide 12: Testing & Validation
add_content_slide("Testing & Validation", [
    "Rigorous 3-Stage Test Strategy:",
    "- ST1: Component-Level Testing (Unit)",
    "- ST2: Sub-System Integration Testing",
    "- ST3: System-Level Validation (End-to-End)",
    "Measured Performance:",
    "FSR -> Valve Latency < 150ms",
    "Pressure Accuracy +/- 1 kPa",
    "Display Refresh Rate ~20 Hz"
])

# Slide 13: Conclusions
add_content_slide("Conclusions", [
    "Achievements & Future Work:",
    "- Pneumatic-assisted ROM restoration verified",
    "- 0-100% grip strength scale implemented",
    "- TENS pain management integrated",
    "- All 9 requirement areas passed tests",
    "Future Work:",
    "Mobile App / PC Dashboard, Wireless Monitoring, Clinical Trials"
])

# Slide 14: Team
add_content_slide("The Team", [
    "Supervised by: Dr. Essam Nabil",
    "Head of Dept: Dr. Ghada El-Banby",
    "Menoufia University",
    "Team of 15 members.",
    "Thank you for your support!"
])

prs.save("Smart_Rehabilitation_Glove.pptx")
print("PowerPoint generated successfully: Smart_Rehabilitation_Glove.pptx")
