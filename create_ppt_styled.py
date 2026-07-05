import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_COLOR = RGBColor(249, 250, 251)
ORANGE = RGBColor(249, 115, 22)
DARK = RGBColor(17, 24, 39)
GRAY = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(229, 231, 235)

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_hero(title, subtitle, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5.5), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(5.5), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY
    
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.0), Inches(1.0), height=Inches(5.5))

def add_slide(title, cards, img_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    # Cards
    top_y = 1.8
    for c in cards:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(top_y), Inches(6.0), Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BORDER
        
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.1)
        tf_card.margin_bottom = Inches(0.1)
        
        if ":" in c:
            parts = c.split(":", 1)
            p1 = tf_card.add_paragraph()
            p1.text = parts[0].strip()
            p1.font.bold = True
            p1.font.size = Pt(16)
            p1.font.color.rgb = DARK
            
            p2 = tf_card.add_paragraph()
            p2.text = parts[1].strip()
            p2.font.size = Pt(14)
            p2.font.color.rgb = GRAY
        else:
            p1 = tf_card.add_paragraph()
            p1.text = c
            p1.font.size = Pt(16)
            p1.font.color.rgb = DARK
            
        top_y += 1.05

    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.8), height=Inches(4.5))
        except Exception as e:
            print("Error adding image:", e)

# Slides
add_hero("Smart Rehabilitation Glove", "A wearable mechatronic device combining TENS therapy, vibration feedback, and real-time force sensing.\n\nGraduation Project 2026\nMenoufia University\nSupervised By: Dr. Essam Nabil", "images/image3.jpeg")

add_slide("Abstract: An Integrated System", [
    "Overview: A smart wearable rehabilitation glove for patients with neurological impairments.",
    "TENS Module: Transcutaneous Electrical Nerve Stimulation for pain modulation.",
    "Vibration Therapy: ERM haptic actuators providing sensory stimulation.",
    "Force Measurement: FSR sensors mapped to 0-100% grip percentage."
], "images/image9.png")

add_slide("Problem Statement", [
    "Limited Range of Motion: Fractures and arthritis restrict joint mobility.",
    "Muscle Atrophy & Weakness: Prolonged immobilization leads to muscle wasting.",
    "Sensory Deficits: Nerve damage alters tactile feedback.",
    "Chronic Pain: Persistent pain acts as a major barrier.",
    "Our Solution: An integrated wearable glove combining multiple therapies."
])

add_slide("Historical Background", [
    "Ancient Era: Primitive splinting and exercises.",
    "18th-19th Century: Surgical pioneers correcting hand deformities.",
    "World Wars: Severe trauma catalyzed advances like FES.",
    "1950s-1977: Hand rehabilitation emerged as a formal specialty.",
    "Today 2026: Smart Rehabilitation Gloves with real-time feedback."
])

add_slide("Objectives (1/2)", [
    "Restoring Range of Motion (ROM): Pneumatic-assisted exercises.",
    "Increasing Strength & Dexterity: Strengthening muscles for fine motor tasks.",
    "Non-Invasive Pain Management: TENS therapy as a safer alternative to drugs."
], "images/image7.png")

add_slide("Objectives (2/2)", [
    "Enhancing Activities of Daily Living: Empowering independent performance.",
    "Improving Coordination & Proprioception: Neuromuscular coordination.",
    "Customized Plans: Individualized protocols tailored to each patient.",
    "Scar Management: Targeted techniques to minimize post-surgical scar tissue."
])

add_slide("System Architecture", [
    "Input & Sensing Layer: Flex sensors + FSR sensors with 12-bit ADC.",
    "Processing & Control Layer: ESP32 dual-core handles logic and PWM.",
    "Actuation & Power Layer: Solenoid valves + MOSFET + Buck converter.",
    "Feedback & Interface Layer: LCD display showing real-time metrics."
], "images/image6.jpeg")

add_slide("Core Components", [
    "ESP32 Microcontroller: Dual-core LX6 processor.",
    "Pneumatic Glove: Inflatable air actuators along each finger.",
    "LCD Display Screen: Shows real-time system data and diagnostics.",
    "MOSFET Module: High-speed switch for valves and motors.",
    "Power Supply Adaptor: Converts 220V AC to DC.",
    "Power Module (LM2596): Step-down converter for stable voltages."
], "images/image9.png")

add_slide("Sensors & Therapy Components", [
    "TENS Cable & Electrodes: Transmits electrical impulses.",
    "Vibration Motor (ERM): Provides tactile feedback.",
    "Force Sensor (FSR): Thin flexible sensor measuring grip pressure.",
    "Solenoid Valves: Electromechanical air valves.",
    "Voltage Regulator 7805: Stable fixed +5V DC output.",
    "Air Pressure Sensor HX710B: 24-bit ADC chip for closed-loop control."
], "images/image14.png")

add_slide("Development Methodology (1/2)", [
    "1. Analysis & Requirements: Defined functional & non-functional requirements.",
    "2. Theoretical Design: System decomposed into 4 layers.",
    "3. Modular Assembly & Unit Testing: Independent component testing."
], "images/image3.jpeg")

add_slide("Development Methodology (2/2)", [
    "4. Full System Integration: Unified firmware managing all tasks.",
    "5. Functional Verification: Tested under real-world clinical conditions.",
    "6. Calibration & Final Tuning: Sensor offset adjusted and valves fine-tuned."
])

add_slide("Testing & Validation", [
    "ST1 Component-Level Testing (Unit): 13 components tested individually.",
    "ST2 Sub-System Integration Testing: Group interactions verified.",
    "ST3 System-Level Validation: Complete clinical scenario looped.",
    "Performance: FSR latency < 150ms, Accuracy +/- 1 kPa."
])

add_slide("Conclusions & Future Work", [
    "Achievements: Pneumatic-assisted ROM restoration verified.",
    "Tracking: 0-100% grip strength scale for quantitative tracking.",
    "Safety: Real-time display with auto pressure shutoff.",
    "Future Enhancements: Mobile App, Wireless Monitoring.",
    "Clinical Trials: Validate therapeutic efficacy across larger groups."
])

add_slide("Work Team", [
    "Supervision: Dr. Essam Nabil & Dr. Ghada El-Banby.",
    "Institution: Menoufia University.",
    "Team Members: 15 dedicated engineers.",
    "Special Thanks: For the invaluable support and dedication."
])

prs.save("Smart_Rehabilitation_Glove_Styled.pptx")
print("Styled PowerPoint created successfully.")
