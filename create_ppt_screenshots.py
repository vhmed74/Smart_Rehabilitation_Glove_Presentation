import os
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

def main():
    html_path = "file:///C:/Users/LOQ/Downloads/Smart_Rehabilitation_Glove_Presentation/index.html"
    scratch_dir = "screenshots"
    os.makedirs(scratch_dir, exist_ok=True)
    
    total_slides = 14
    
    print("Starting Playwright...")
    with sync_playwright() as p:
        # Launch Chromium (headless)
        browser = p.chromium.launch(headless=True)
        # 16:9 ratio, standard FHD
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        
        print("Loading HTML presentation...")
        page.goto(html_path)
        
        # Give it a moment to load Swiper and GSAP
        page.wait_for_timeout(2000)
        
        # We can hide the pagination and counter if we want a clean look, 
        # but leaving them gives the "presentation" feel.
        # Let's hide the scroll hint so it doesn't blink in screenshots
        page.evaluate("document.querySelector('.scroll-h').style.display = 'none';")
        
        for i in range(total_slides):
            print(f"Capturing slide {i+1}/{total_slides}...")
            # Move swiper instantly
            page.evaluate(f"swiper.slideTo({i}, 0);")
            
            # Wait for GSAP animations to finish. Most are 0.7s, so 1.5s is safe
            page.wait_for_timeout(1500)
            
            # Take screenshot
            shot_path = os.path.join(scratch_dir, f"slide_{i}.png")
            page.screenshot(path=shot_path)
            
        browser.close()
        
    print("Screenshots complete. Building PowerPoint...")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i in range(total_slides):
        slide_layout = prs.slide_layouts[6] # blank
        slide = prs.slides.add_slide(slide_layout)
        shot_path = os.path.join(scratch_dir, f"slide_{i}.png")
        if os.path.exists(shot_path):
            slide.shapes.add_picture(shot_path, 0, 0, width=Inches(13.333))
            
    output_file = "Smart_Rehabilitation_Glove_Original_Style.pptx"
    prs.save(output_file)
    print(f"Success! Saved as {output_file}")

if __name__ == "__main__":
    main()
